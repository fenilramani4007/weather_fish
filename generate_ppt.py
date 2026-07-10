"""
WEATHER-FISH - Simple Overview Presentation
Run: py -3.11 generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
BG    = RGBColor(0x07, 0x09, 0x0f)
PANEL = RGBColor(0x0f, 0x18, 0x28)
GOLD  = RGBColor(0xd4, 0xb4, 0x5a)
TEXT  = RGBColor(0xea, 0xf1, 0xf8)
GRAY  = RGBColor(0x9e, 0xb8, 0xd0)
GRN   = RGBColor(0x4a, 0xde, 0x80)
BLUE  = RGBColor(0x60, 0xa5, 0xfa)
DARK  = RGBColor(0x15, 0x22, 0x35)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, fill, border=None):
    sh = s.shapes.add_shape(1, x, y, w, h)
    from pptx.oxml.ns import qn
    pg = sh.element.spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "rect")
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def text(s, t, x, y, w, h, sz=14, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = t
    r.font.name = "Calibri"
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def slide_header(s, title, num):
    rect(s, 0, 0, Inches(0.08), H, GOLD)
    rect(s, Inches(0.08), 0, W, Inches(1.0), PANEL)
    rect(s, Inches(0.08), Inches(1.0), W, Inches(0.04), GOLD)
    text(s, title, Inches(0.5), Inches(0.18), W - Inches(2), Inches(0.7),
         sz=28, color=GOLD, bold=True)
    text(s, f"{num} / 8", W - Inches(1.1), Inches(0.3), Inches(0.9), Inches(0.35),
         sz=10, color=GRAY, align=PP_ALIGN.RIGHT)
    return Inches(1.2)


# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s1 = slide()

# Gold top & bottom bars
rect(s1, 0, 0, W, Inches(0.12), GOLD)
rect(s1, 0, H - Inches(0.12), W, Inches(0.12), GOLD)

# Center panel
rect(s1, Inches(1.5), Inches(1.4), W - Inches(3.0), H - Inches(2.8), PANEL)

text(s1, "WEATHER-FISH",
     Inches(1.5), Inches(1.6), W - Inches(3.0), Inches(1.6),
     sz=52, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font="Calibri")

text(s1, "AI-Driven Adaptive Weather Narration System",
     Inches(1.5), Inches(3.0), W - Inches(3.0), Inches(0.6),
     sz=18, color=TEXT, align=PP_ALIGN.CENTER)

text(s1, "OTH Amberg-Weiden  |  Summer Semester 2026",
     Inches(1.5), Inches(3.72), W - Inches(3.0), Inches(0.45),
     sz=13, color=GRAY, align=PP_ALIGN.CENTER)

text(s1, "Fenil Ramani   Prasad Rajyguru   Shubham Kushwaha",
     Inches(1.5), Inches(4.28), W - Inches(3.0), Inches(0.42),
     sz=13, color=GRAY, align=PP_ALIGN.CENTER)

text(s1, "FastAPI  |  Google Gemini AI  |  React 18  |  MongoDB  |  Edge-TTS  |  JWT",
     Inches(1.5), Inches(5.0), W - Inches(3.0), Inches(0.38),
     sz=11, color=GOLD, align=PP_ALIGN.CENTER, italic=True)


# ── SLIDE 2: What is WEATHER-FISH? ───────────────────────────────────────────
s2 = slide()
cy = slide_header(s2, "What is WEATHER-FISH?", 2)

text(s2,
     "Transforms a German postal code into intelligent, character-styled audio weather reports.",
     Inches(0.5), cy, W - Inches(1.0), Inches(0.6),
     sz=17, color=TEXT)

cy += Inches(0.75)

points = [
    ("Three AI Presenters",  "Fisch (playful) · Merkel (formal) · Haftbefehl (street-rap style)"),
    ("Bilingual Output",     "All reports generated in German AND English simultaneously"),
    ("AI-Powered Narration", "Google Gemini writes context-aware, persona-specific scripts"),
    ("Neural Audio",         "Microsoft Edge-TTS converts text to natural-sounding MP3 voices"),
    ("Fully Deployed",       "Live on HuggingFace Spaces — publicly accessible"),
]

for title, desc in points:
    rect(s2, Inches(0.5), cy, W - Inches(1.0), Inches(0.7), PANEL, GOLD)
    rect(s2, Inches(0.5), cy, Inches(0.07), Inches(0.7), GOLD)
    text(s2, title, Inches(0.75), cy + Inches(0.08), Inches(3.0), Inches(0.28),
         sz=12.5, color=GOLD, bold=True)
    text(s2, desc, Inches(3.9), cy + Inches(0.08), W - Inches(4.6), Inches(0.28),
         sz=12.5, color=TEXT)
    cy += Inches(0.78)


# ── SLIDE 3: How It Works ────────────────────────────────────────────────────
s3 = slide()
cy = slide_header(s3, "How It Works", 3)

steps = [
    ("1", "Postal Code"),
    ("2", "Geocoding"),
    ("3", "Weather Data"),
    ("4", "AI Context"),
    ("5", "Gemini AI"),
    ("6", "Edge-TTS"),
    ("7", "Report + Audio"),
]

bw = Inches(1.6)
bh = Inches(1.8)
gap = Inches(0.08)
total = len(steps) * bw + (len(steps) - 1) * gap
sx = (W - total) / 2

for i, (num, label) in enumerate(steps):
    bx = sx + i * (bw + gap)
    by = cy + Inches(0.3)
    rect(s3, bx, by, bw, bh, PANEL, GOLD)
    text(s3, num, bx, by + Inches(0.15), bw, Inches(0.7),
         sz=30, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    text(s3, label, bx + Inches(0.08), by + Inches(0.85), bw - Inches(0.16), Inches(0.65),
         sz=11.5, color=TEXT, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        text(s3, "->", bx + bw + Inches(0.01), by + Inches(0.65),
             gap + Inches(0.06), Inches(0.4),
             sz=14, color=GRAY, align=PP_ALIGN.CENTER)

text(s3,
     "6 parallel tasks (ThreadPoolExecutor)  —  Generation time reduced from ~90s to ~15s",
     Inches(0.5), cy + Inches(2.5), W - Inches(1.0), Inches(0.4),
     sz=13, color=GRN, align=PP_ALIGN.CENTER, italic=True)

text(s3,
     "Reports cached in MongoDB for 20 minutes — skips AI entirely on repeat requests",
     Inches(0.5), cy + Inches(2.95), W - Inches(1.0), Inches(0.4),
     sz=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ── SLIDE 4: Key Features ─────────────────────────────────────────────────────
s4 = slide()
cy = slide_header(s4, "Key Features", 4)

features = [
    "6-model Gemini fallback chain (2.5-flash to 1.5-flash-8b)",
    "Template fallback — always works even without API quota",
    "Bilingual: German + English simultaneously (3 x 2 = 6 reports)",
    "Neural TTS voices — 6 MP3 files per generation",
    "Context Engine — severity, time-of-day, adaptive tone",
    "MongoDB caching — 20-min TTL, 90-day weather history",
    "APScheduler — auto-regenerates reports every hour",
    "JWT Authentication — register, login, user profiles",
    "Activity tracking — logs daily searches by location and hobby",
    "Mobile-responsive — hamburger nav drawer at 640px",
]

col_w = (W - Inches(1.1)) / 2 - Inches(0.1)
for i, feat in enumerate(features):
    col = i % 2
    row = i // 2
    fx = Inches(0.5) + col * (col_w + Inches(0.2))
    fy = cy + Inches(0.05) + row * Inches(0.72)
    text(s4, "  +  " + feat, fx, fy, col_w, Inches(0.5),
         sz=12.5, color=TEXT)
    rect(s4, fx, fy + Inches(0.58), col_w, Inches(0.01), DARK)


# ── SLIDE 5: Tech Stack ───────────────────────────────────────────────────────
s5 = slide()
cy = slide_header(s5, "Tech Stack", 5)

stack = [
    ("Backend",     "FastAPI + Python 3.11 + Uvicorn",        GOLD),
    ("AI",          "Google Gemini (google-genai 2.8.0)",      GOLD),
    ("Frontend",    "React 18 + TypeScript + Vite 6",          BLUE),
    ("Design",      "German Bauhaus dark theme (gold/black)",  BLUE),
    ("Database",    "MongoDB Atlas M0 + pymongo 4.7.3",        GRN),
    ("TTS",         "Microsoft Edge-TTS (Neural voices)",      GRN),
    ("Auth",        "JWT (PyJWT) + bcrypt (direct)",           GRAY),
    ("Hosting",     "HuggingFace Spaces (live deployment)",    GRAY),
]

col_w5 = (W - Inches(1.1)) / 2 - Inches(0.1)
for i, (cat, val, col) in enumerate(stack):
    c = i % 2
    r = i // 2
    bx = Inches(0.5) + c * (col_w5 + Inches(0.2))
    by = cy + Inches(0.08) + r * Inches(0.82)
    rect(s5, bx, by, col_w5, Inches(0.72), PANEL, col)
    rect(s5, bx, by, Inches(0.07), Inches(0.72), col)
    text(s5, cat, bx + Inches(0.18), by + Inches(0.1), Inches(1.6), Inches(0.28),
         sz=12, color=col, bold=True)
    text(s5, val, bx + Inches(1.9), by + Inches(0.1), col_w5 - Inches(2.1), Inches(0.28),
         sz=12, color=TEXT)


# ── SLIDE 6: Team ─────────────────────────────────────────────────────────────
s6 = slide()
cy = slide_header(s6, "Team", 6)

members = [
    {
        "name": "Fenil Ramani",
        "role": "Project Lead",
        "color": GOLD,
        "items": [
            "AI engine + prompt engineering",
            "Context Engine (severity, tone)",
            "Edge-TTS integration (DE + EN)",
            "JWT auth + user system",
            "MongoDB schema + caching",
            "HuggingFace deployment",
        ],
    },
    {
        "name": "Prasad Rajyguru",
        "role": "Frontend Engineer",
        "color": BLUE,
        "items": [
            "React 18 dashboard (Bauhaus theme)",
            "3 mascot presenter UI cards",
            "Audio player + wave animation",
            "Mobile responsive navigation",
            "Reports page DE/EN toggle",
        ],
    },
    {
        "name": "Shubham Kushwaha",
        "role": "Backend Engineer",
        "color": GRN,
        "items": [
            "FastAPI REST API (15 endpoints)",
            "OpenWeather + Geocoding pipeline",
            "APScheduler hourly auto-generation",
            "MongoDB CRUD (pymongo)",
        ],
    },
]

mw = (W - Inches(1.0)) / 3 - Inches(0.12)
mh = H - cy - Inches(0.4)

for i, m in enumerate(members):
    mx = Inches(0.5) + i * (mw + Inches(0.16))
    rect(s6, mx, cy, mw, mh, PANEL)
    rect(s6, mx, cy, mw, Inches(0.55), DARK)
    rect(s6, mx, cy, Inches(0.07), Inches(0.55), m["color"])
    text(s6, m["name"], mx + Inches(0.18), cy + Inches(0.06),
         mw - Inches(0.22), Inches(0.28), sz=13, color=m["color"], bold=True)
    text(s6, m["role"], mx + Inches(0.18), cy + Inches(0.32),
         mw - Inches(0.22), Inches(0.22), sz=10, color=GRAY, italic=True)
    iy = cy + Inches(0.65)
    for item in m["items"]:
        text(s6, "  -  " + item, mx + Inches(0.1), iy,
             mw - Inches(0.2), Inches(0.3), sz=11, color=TEXT)
        iy += Inches(0.35)


# ── SLIDE 7: Project Status ───────────────────────────────────────────────────
s7 = slide()
cy = slide_header(s7, "Project Status", 7)

chips = [
    ("LIVE",    "HuggingFace Spaces", GRN),
    ("100%",    "Core Pipeline Done", GRN),
    ("~15s",    "Generation Time",    GOLD),
    ("JWT",     "Auth Active",        BLUE),
]
cw7 = (W - Inches(1.0)) / 4 - Inches(0.12)
for i, (val, label, col) in enumerate(chips):
    cx7 = Inches(0.5) + i * (cw7 + Inches(0.14))
    rect(s7, cx7, cy + Inches(0.1), cw7, Inches(1.1), PANEL, col)
    text(s7, val, cx7, cy + Inches(0.15), cw7, Inches(0.55),
         sz=24, color=col, bold=True, align=PP_ALIGN.CENTER)
    text(s7, label, cx7, cy + Inches(0.68), cw7, Inches(0.35),
         sz=11, color=GRAY, align=PP_ALIGN.CENTER)

cy += Inches(1.4)

hw = (W - Inches(1.1)) / 2 - Inches(0.1)

text(s7, "Completed", Inches(0.5), cy, hw, Inches(0.35),
     sz=14, color=GRN, bold=True)
text(s7, "Next Steps", Inches(0.7) + hw, cy, hw, Inches(0.35),
     sz=14, color=GOLD, bold=True)
cy += Inches(0.42)

done = [
    "End-to-end AI pipeline",
    "Bilingual DE + EN reports",
    "6-model Gemini fallback",
    "Parallel generation (6x faster)",
    "JWT auth + user profiles",
    "Mobile responsive UI",
    "HuggingFace deployment",
]
next_steps = [
    "Final QA and testing - July 2026",
    "Documentation completion",
    "UI accessibility polish",
    "Final submission & presentation",
    "Future: Historical trend charts",
    "Future: Weather alert notifications",
]

for i in range(max(len(done), len(next_steps))):
    if i < len(done):
        text(s7, "  +  " + done[i], Inches(0.5), cy + i * Inches(0.44),
             hw, Inches(0.38), sz=12, color=TEXT)
    if i < len(next_steps):
        text(s7, "  ->  " + next_steps[i], Inches(0.7) + hw, cy + i * Inches(0.44),
             hw, Inches(0.38), sz=12, color=TEXT)


# ── SLIDE 8: Thank You ────────────────────────────────────────────────────────
s8 = slide()

rect(s8, 0, 0, W, Inches(0.12), GOLD)
rect(s8, 0, H - Inches(0.12), W, Inches(0.12), GOLD)
rect(s8, Inches(1.5), Inches(1.8), W - Inches(3.0), H - Inches(3.6), PANEL)

text(s8, "Thank You",
     Inches(1.5), Inches(2.0), W - Inches(3.0), Inches(1.1),
     sz=48, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

text(s8, "WEATHER-FISH  |  OTH Amberg-Weiden  |  Summer 2026",
     Inches(1.5), Inches(3.3), W - Inches(3.0), Inches(0.5),
     sz=14, color=TEXT, align=PP_ALIGN.CENTER)

text(s8, "Fenil Ramani   |   Prasad Rajyguru   |   Shubham Kushwaha",
     Inches(1.5), Inches(3.9), W - Inches(3.0), Inches(0.45),
     sz=13, color=GRAY, align=PP_ALIGN.CENTER)

text(s8, "AI Tool Disclosure: Claude (Anthropic) used for code implementation and debugging. "
         "All code reviewed and integrated by the team.",
     Inches(1.5), Inches(5.7), W - Inches(3.0), Inches(0.5),
     sz=9, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "WEATHER-FISH_Presentation.pptx"
prs.save(OUT)
print(f"Saved -> {OUT}  ({len(prs.slides)} slides)")
